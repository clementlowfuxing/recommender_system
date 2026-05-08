"""
Generate the CDWC Diagnostic Engine — High-Level Architecture & Scoring Flow deck.
Scope aligned with .kiro/specs/recommendation-engine-redesign/requirements.md:
- Engine stops at diagnosis. Action triage (training / exposure / hiring)
  is OUT OF SCOPE.
- Per-cell TPCP competency fields (b*, k*, p*, e*) are OUT OF SCOPE.
  Only block-level rollup percentages are consumed.
- 5 per-person health scores; 2 parallel analyses (team aggregate +
  individual drill-down).

Output: presentation_ppt/CDWC_Architecture_Scoring_Flow.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "presentation_ppt"
OUTPUT_DIR.mkdir(exist_ok=True)
OUTPUT_FILE = OUTPUT_DIR / "CDWC_Architecture_Scoring_Flow.pptx"

# Palette
DARK_BLUE = RGBColor(0x1E, 0x29, 0x3B)
PRIMARY_BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
PALE_BLUE = RGBColor(0xF0, 0xF7, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
GREEN = RGBColor(0x05, 0x96, 0x69)
PALE_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
RED = RGBColor(0xDC, 0x26, 0x26)
PALE_RED = RGBColor(0xFE, 0xE2, 0xE2)
YELLOW_BG = RGBColor(0xFE, 0xF9, 0xC3)
AMBER = RGBColor(0xCA, 0x8A, 0x04)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PALE_PURPLE = RGBColor(0xF3, 0xE8, 0xFF)
TEAL = RGBColor(0x0D, 0x94, 0x88)
PALE_TEAL = RGBColor(0xCC, 0xFB, 0xF1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1.25)
    else:
        s.line.fill.background()
    return s


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK_BLUE, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def add_bullets(slide, left, top, width, height, items, size=14, color=DARK_BLUE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, weight=Pt(2)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = weight
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = c.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return c


def add_header(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.65),
             title, size=26, bold=True, color=DARK_BLUE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.4),
                 subtitle, size=13, color=GRAY)
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(0.6), Inches(0.05), PRIMARY_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
         "CDWC Diagnostic Engine",
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         "High-Level Architecture & Scoring Flow",
         size=24, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.6),
         "How the Engine Turns SMA & TPCP Data into Diagnostic JSON",
         size=15, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
         "Diagnosis only  ·  Block-level rollups only  ·  Action triage out of scope",
         size=13, bold=True, color=RGBColor(0xFB, 0xBF, 0x24), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — End-to-End Architecture
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "End-to-End Architecture",
           "From HOD/GM question → to data sources → to scoring → to diagnostic JSON")

layers = [
    ("① PRESENTATION LAYER   ·   HOD/GM asks a question", PALE_PURPLE, PURPLE),
    ("② CONVERSATION LAYER   ·   NL parser + narrative formatter", PALE_BLUE, PRIMARY_BLUE),
    ("③ API LAYER   ·   FastAPI endpoints route the request", LIGHT_BLUE, PRIMARY_BLUE),
    ("④ ENGINE LAYER   ·   join → scope → score → analyze   (PURE JSON OUT — DIAGNOSIS ONLY)", PALE_GREEN, GREEN),
    ("⑤ DATA LAYER   ·   DUMMY_SMA_HRIS_sample.json  +  DUMMY_TPCP_Assessment_sample.json", LIGHT_GRAY, GRAY),
]

layer_h = Inches(0.9)
y_start = Inches(1.6)
for i, (lbl, fill, border) in enumerate(layers):
    y = y_start + i * (layer_h + Inches(0.15))
    add_shape(slide, Inches(0.5), y, Inches(12.3), layer_h, fill, border)
    add_text(slide, Inches(0.8), y + Inches(0.25), Inches(11.7), Inches(0.45),
             lbl, size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)

for i in range(4):
    y = y_start + (i + 1) * (layer_h + Inches(0.15)) - Inches(0.14)
    add_text(slide, Inches(6.2), y, Inches(1), Inches(0.25),
             "↕", size=14, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
         "Engine layer never formats prose. It emits structured diagnostic JSON; action recommendations are handled by a separate future layer.",
         size=12, bold=True, color=PRIMARY_BLUE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Layer zoom
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Data Layer",
           "Two source files, full-outer-join on dummy_id → unified employee record")

# HRIS box
add_shape(slide, Inches(0.5), Inches(1.75), Inches(5.2), Inches(3.2), LIGHT_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(0.7), Inches(1.9), Inches(4.8), Inches(0.45),
         "DUMMY_SMA_HRIS_sample.json", size=14, bold=True, color=PRIMARY_BLUE)
add_text(slide, Inches(0.7), Inches(2.35), Inches(4.8), Inches(0.4),
         "9,455 employees  (1 row per person)", size=11, color=GRAY)
hris_fields = [
    "• dummy_id, employee_name, superior_name",
    "• skill_group, job_grade, salary_grade",
    "• year_in_sg",
    "• sma_completion_status",
    "• overall_cbs_pct",
    "• technical_cbs_pct",
    "• leadership_cbs_pct",
]
add_bullets(slide, Inches(0.8), Inches(2.8), Inches(4.7), Inches(2), hris_fields, size=12)

# TPCP box
add_shape(slide, Inches(7.6), Inches(1.75), Inches(5.2), Inches(3.2), PALE_PURPLE, PURPLE)
add_text(slide, Inches(7.8), Inches(1.9), Inches(4.8), Inches(0.45),
         "DUMMY_TPCP_Assessment_sample.json", size=14, bold=True, color=PURPLE)
add_text(slide, Inches(7.8), Inches(2.35), Inches(4.8), Inches(0.4),
         "1,018 assessments  (1 row per person)", size=11, color=GRAY)
tpcp_fields = [
    "• dummy_id, employee_name",
    "• business, opu, division, skill_group",
    "• discipline, sub_discipline",
    "• assessment_level, qualified_level, passed_tpcp",
    "• base_pct, key_pct, pacing_pct, emerging_pct, cti_met_pct",
    "• assessment_date   (per-cell b/k/p/e fields NOT consumed)",
]
add_bullets(slide, Inches(7.8), Inches(2.8), Inches(4.7), Inches(2), tpcp_fields, size=11)

# Join arrows and box
add_arrow(slide, Inches(5.7), Inches(3.3), Inches(6.4), Inches(3.3), PRIMARY_BLUE, Pt(2.5))
add_arrow(slide, Inches(7.6), Inches(3.3), Inches(6.9), Inches(3.3), PURPLE, Pt(2.5))
add_shape(slide, Inches(5.7), Inches(2.9), Inches(1.9), Inches(0.8), YELLOW_BG, AMBER)
add_text(slide, Inches(5.7), Inches(3.0), Inches(1.9), Inches(0.5),
         "FULL JOIN", size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_text(slide, Inches(5.7), Inches(3.35), Inches(1.9), Inches(0.35),
         "on dummy_id", size=10, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Unified record
add_shape(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(1.7), PALE_GREEN, GREEN)
add_text(slide, Inches(1.7), Inches(5.4), Inches(10), Inches(0.4),
         "Unified Employee Record", size=14, bold=True, color=GREEN)
unified_desc = [
    "has_hris + has_tpcp flags preserved — BOTH (full signal)  ·  HRIS-only (CBS only)  ·  TPCP-only (edge case)",
    "Each record carries identity + org hierarchy + CBS block + TPCP rollups (base/key/pacing/emerging %, cti_met %).",
]
add_bullets(slide, Inches(1.7), Inches(5.85), Inches(10), Inches(1.1), unified_desc, size=12)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — The 4-stage pipeline (diagnosis only)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Engine Pipeline — 4 Stages (Diagnosis Only)",
           "Every query flows through these stages; outputs are pure diagnostic JSON")

stages = [
    {"n": "1", "t": "Load & Join", "c": PALE_TEAL, "b": TEAL,
     "d": "Ingest HRIS + TPCP.\nFull-outer-join on\ndummy_id. Flag each\nrecord has_hris / has_tpcp."},
    {"n": "2", "t": "Scope Filter", "c": PALE_RED, "b": RED,
     "d": "Apply mandatory org filter:\nbusiness · opu · division.\nResulting team = scoped\npool. Signal, not gate."},
    {"n": "3", "t": "Per-Person Scoring", "c": LIGHT_BLUE, "b": PRIMARY_BLUE,
     "d": "Compute 5 health scores\nfor every scoped employee.\nAll scores are 0–1.\nFallbacks for missing data."},
    {"n": "4", "t": "Analysis Layer", "c": YELLOW_BG, "b": AMBER,
     "d": "Run 2 parallel analyses:\n• Team aggregate rollup\n• Individual drill-down\n(flags + lists)"},
]

box_w = Inches(2.9)
box_h = Inches(3.2)
y = Inches(1.7)
start_x = Inches(0.5)
gap = Inches(0.25)

for i, s in enumerate(stages):
    x = start_x + i * (box_w + gap)
    add_shape(slide, x, y, box_w, box_h, s["c"], s["b"])
    add_shape(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5), s["b"], s["b"], shape=MSO_SHAPE.OVAL)
    add_text(slide, x + Inches(0.2), y + Inches(0.28), Inches(0.5), Inches(0.4),
             s["n"], size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.85), y + Inches(0.28), box_w - Inches(1), Inches(0.45),
             s["t"], size=16, bold=True, color=DARK_BLUE)
    add_text(slide, x + Inches(0.25), y + Inches(1), box_w - Inches(0.5), box_h - Inches(1.1),
             s["d"], size=12, color=DARK_BLUE)
    if i < 3:
        add_text(slide, x + box_w, y + Inches(1.4), gap, Inches(0.4),
                 "→", size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Question routing below
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(0.7), Inches(5.3), Inches(12), Inches(0.4),
         "Which stage answers which question?", size=14, bold=True, color=PRIMARY_BLUE)
route = [
    "• Stage 1–2  →  'Who is in my scope?'  'Who has HRIS but no TPCP?'  'Team size after filter?'",
    "• Stage 3    →  'Who has not completed TPCP?'  'Individual scorecards'  'What is their combined score?'",
    "• Stage 4    →  'Team TPCP performance?'  'Which block is weakest?'  'Who is promotable?'  'Who are the low performers?'",
]
add_bullets(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(1.3), route, size=12)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Scoring engine: inputs → 5 health scores → outputs
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Per-Person Scoring — Inputs to 5 Health Scores",
           "Each employee becomes a 5-dimensional health vector")

# Left column — data inputs
add_shape(slide, Inches(0.4), Inches(1.7), Inches(3.2), Inches(5.3), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(0.6), Inches(1.8), Inches(2.8), Inches(0.4),
         "INPUT FIELDS", size=13, bold=True, color=PRIMARY_BLUE, align=PP_ALIGN.CENTER)
inputs_block = [
    "FROM HRIS:",
    "  • overall_cbs_pct",
    "  • technical_cbs_pct",
    "  • leadership_cbs_pct",
    "  • job_grade",
    "  • year_in_sg",
    "",
    "FROM TPCP:",
    "  • cti_met_pct",
    "  • base_pct",
    "  • key_pct",
    "  • pacing_pct",
    "  • emerging_pct",
    "  • passed_tpcp",
    "  • assessment_level",
    "",
    "  (per-cell b/k/p/e",
    "   fields excluded)",
]
add_bullets(slide, Inches(0.6), Inches(2.3), Inches(2.8), Inches(4.7), inputs_block, size=11)

# Middle — 5 score boxes (3+2 layout centered)
score_titles = [
    ("1. TPCP Health",        "cti_met_pct, base, key,\npacing, emerging, pass", LIGHT_BLUE, PRIMARY_BLUE),
    ("2. CBS Health",         "technical_cbs +\noverall_cbs +\nleadership_cbs", PALE_PURPLE, PURPLE),
    ("3. Combined Competency","0.55 × TPCP\n+ 0.45 × CBS", PALE_GREEN, GREEN),
    ("4. Grade Fit",          "job_grade ↔\nassessment_level\nalignment", YELLOW_BG, AMBER),
    ("5. Tenure Readiness",   "year_in_sg /\ngrade_threshold", PALE_TEAL, TEAL),
]
col_x = Inches(3.9)
card_w = Inches(3.0)
card_h = Inches(1.6)
# Arrange: row 0 has 3 cards, row 1 has 2 cards (centered)
for i, (t, d, c, b) in enumerate(score_titles):
    if i < 3:
        row = 0
        col = i
        x = col_x + col * (card_w + Inches(0.15))
    else:
        row = 1
        col = i - 3
        # center 2 cards under 3 cards above
        left_shift = (card_w + Inches(0.15)) * 0.5
        x = col_x + left_shift + col * (card_w + Inches(0.15))
    y = Inches(1.9) + row * (card_h + Inches(0.3))
    add_shape(slide, x, y, card_w, card_h, c, b)
    add_text(slide, x + Inches(0.15), y + Inches(0.1), card_w - Inches(0.3), Inches(0.45),
             t, size=13, bold=True, color=DARK_BLUE)
    add_text(slide, x + Inches(0.15), y + Inches(0.55), card_w - Inches(0.3), card_h - Inches(0.65),
             d, size=11, color=DARK_BLUE)

# Right — output vector
right_x = col_x + 3 * (card_w + Inches(0.15)) + Inches(0.05)
add_shape(slide, right_x, Inches(1.7), Inches(2.6), Inches(5.3), DARK_BLUE, DARK_BLUE)
add_text(slide, right_x + Inches(0.2), Inches(1.85), Inches(2.2), Inches(0.4),
         "OUTPUT", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb = slide.shapes.add_textbox(right_x + Inches(0.2), Inches(2.3), Inches(2.3), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
out_text = """{
  "dummy_id": 23267,
  "tpcp_health": 0.78,
  "cbs_health": 0.64,
  "combined_competency":
    0.72,
  "grade_fit": 1.0,
  "tenure_readiness":
    0.85
}"""
p = tf.paragraphs[0]
p.text = out_text
p.font.size = Pt(10)
p.font.name = "Consolas"
p.font.color.rgb = WHITE


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Deep dive: how each score is computed (5 scores)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "How Each Score Is Computed",
           "Every formula is deterministic, traceable, and tunable via config")

score_details = [
    ("1. TPCP Health", LIGHT_BLUE, PRIMARY_BLUE,
     "tpcp_health =\n   0.6 × cti_met_pct\n + 0.4 × avg(base_pct, key_pct,\n            pacing_pct, emerging_pct)\n + 0.1 if passed_tpcp = 'Y'\n (capped at 1.0)",
     "→ 0 if no TPCP record"),
    ("2. CBS Health", PALE_PURPLE, PURPLE,
     "cbs_health =\n   0.5 × technical_cbs_pct\n + 0.3 × overall_cbs_pct\n + 0.2 × leadership_cbs_pct\n (weights rebalance if a\n  dimension is missing)",
     "→ available to all HRIS rows"),
    ("3. Combined Competency", PALE_GREEN, GREEN,
     "combined_competency =\n   0.55 × tpcp_health\n + 0.45 × cbs_health\n\n (if one input is 0 →\n  combined = the other)",
     "→ primary ranking score"),
    ("4. Grade Fit", YELLOW_BG, AMBER,
     "Expected level by grade:\n  P1–P3 → Staff\n  P4–P6 → Principal\n  P7+   → Custodian\n\nExact=1.0  Over=0.5  Under=0.0",
     "→ flags grade/level misalignment"),
    ("5. Tenure Readiness", PALE_TEAL, TEAL,
     "Grade tenure thresholds:\n  P1–P3 → 3 yrs\n  P4–P6 → 4 yrs\n  P7+   → 5 yrs\n\nscore = min(year_in_sg /\n            threshold, 1.0)",
     "→ gates promotability"),
]
card_w = Inches(4.1)
card_h = Inches(2.6)
gx = Inches(0.4)
gy = Inches(1.6)
# Row 0 = 3 cards, Row 1 = 2 cards centered
for i, (t, c, b, formula, note) in enumerate(score_details):
    if i < 3:
        row = 0
        col = i
        x = gx + col * (card_w + Inches(0.1))
    else:
        row = 1
        col = i - 3
        left_shift = (card_w + Inches(0.1)) * 0.5
        x = gx + left_shift + col * (card_w + Inches(0.1))
    y = gy + row * (card_h + Inches(0.2))
    add_shape(slide, x, y, card_w, card_h, c, b)
    add_text(slide, x + Inches(0.15), y + Inches(0.1), card_w - Inches(0.3), Inches(0.4),
             t, size=14, bold=True, color=DARK_BLUE)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), card_w - Inches(0.3), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = formula
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = DARK_BLUE
    add_text(slide, x + Inches(0.15), y + card_h - Inches(0.45), card_w - Inches(0.3), Inches(0.35),
             note, size=10, bold=True, color=b)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Worked example (single employee, 5 scores)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Worked Example — Single Employee",
           "Tracing one record through the scoring engine")

# Input data panel
add_shape(slide, Inches(0.4), Inches(1.65), Inches(4.0), Inches(5.4), LIGHT_GRAY, GRAY)
add_text(slide, Inches(0.55), Inches(1.75), Inches(3.7), Inches(0.4),
         "INPUT  (from joined record)", size=13, bold=True, color=DARK_BLUE)
tb = slide.shapes.add_textbox(Inches(0.55), Inches(2.2), Inches(3.7), Inches(4.7))
tf = tb.text_frame
tf.word_wrap = True
input_text = """dummy_id            : 23267
job_grade           : P4
year_in_sg          : 4.2

From HRIS
  overall_cbs_pct   : 0.65
  technical_cbs_pct : 0.70
  leadership_cbs_pct: 0.55

From TPCP
  assessment_level  : Principal
  passed_tpcp       : Y
  cti_met_pct       : 0.83
  base_pct          : 0.75
  key_pct           : 0.70
  pacing_pct        : 0.60
  emerging_pct      : 0.80

(per-cell fields ignored)"""
p = tf.paragraphs[0]
p.text = input_text
p.font.size = Pt(10.5)
p.font.name = "Consolas"
p.font.color.rgb = DARK_BLUE

add_text(slide, Inches(4.45), Inches(4.0), Inches(0.5), Inches(0.6),
         "→", size=32, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Computation panel
add_shape(slide, Inches(5.0), Inches(1.65), Inches(4.5), Inches(5.4), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(5.15), Inches(1.75), Inches(4.2), Inches(0.4),
         "COMPUTATION", size=13, bold=True, color=PRIMARY_BLUE)
tb = slide.shapes.add_textbox(Inches(5.15), Inches(2.2), Inches(4.2), Inches(4.7))
tf = tb.text_frame
tf.word_wrap = True
calc_text = """tpcp_avg = mean(0.75, 0.70,
                0.60, 0.80) = 0.7125
tpcp_health = 0.6 × 0.83
            + 0.4 × 0.7125
            + 0.1 (passed)
            = 0.49 + 0.285 + 0.1
            = 0.87  (capped 1.0)

cbs_health = 0.5 × 0.70
           + 0.3 × 0.65
           + 0.2 × 0.55
           = 0.35+0.195+0.11 = 0.65

combined_competency
         = 0.55 × 0.87
         + 0.45 × 0.65
         = 0.48 + 0.29 = 0.77

grade_fit  : P4 → Principal ✓  = 1.0
tenure     : 4.2/4.0 capped    = 1.0"""
p = tf.paragraphs[0]
p.text = calc_text
p.font.size = Pt(10)
p.font.name = "Consolas"
p.font.color.rgb = DARK_BLUE

add_text(slide, Inches(9.55), Inches(4.0), Inches(0.5), Inches(0.6),
         "→", size=32, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Output panel
add_shape(slide, Inches(10.1), Inches(1.65), Inches(2.8), Inches(5.4), DARK_BLUE, DARK_BLUE)
add_text(slide, Inches(10.25), Inches(1.75), Inches(2.5), Inches(0.4),
         "OUTPUT (JSON)", size=13, bold=True, color=WHITE)
tb = slide.shapes.add_textbox(Inches(10.25), Inches(2.2), Inches(2.5), Inches(4.7))
tf = tb.text_frame
tf.word_wrap = True
out = """{
  "tpcp_health": 0.87,
  "cbs_health": 0.65,
  "combined_competency":
    0.77,
  "grade_fit": 1.0,
  "tenure_readiness":
    1.0,
  "promotable": true
}"""
p = tf.paragraphs[0]
p.text = out
p.font.size = Pt(11)
p.font.name = "Consolas"
p.font.color.rgb = WHITE


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Aggregation: per-person scores → team insights (2 branches)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "From Per-Person Scores to Team Insights",
           "Stage 4 runs 2 analyses in parallel over the scored population — engine stops here")

# Center input source
add_shape(slide, Inches(5.1), Inches(1.7), Inches(3.2), Inches(1.1), DARK_BLUE, DARK_BLUE)
add_text(slide, Inches(5.2), Inches(1.8), Inches(3), Inches(0.5),
         "Scored Population", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(5.2), Inches(2.25), Inches(3), Inches(0.4),
         "N employees × 5 scores each",
         size=11, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)

# 2 branch arrows
add_arrow(slide, Inches(6.7), Inches(2.8), Inches(3.0), Inches(3.5), GRAY, Pt(2))
add_arrow(slide, Inches(6.7), Inches(2.8), Inches(10.4), Inches(3.5), GRAY, Pt(2))

branches = [
    {
        "title": "A. Team Aggregate", "x": Inches(1.1), "color": PALE_PURPLE, "border": PURPLE,
        "body": [
            "Headline team statistics:",
            "",
            "• team_size, assessed_count,",
            "   not_assessed_count",
            "• pass_rate",
            "• level_distribution",
            "   (Staff, Principal, Custodian,",
            "   Not TP)",
            "• avg_combined_competency",
            "• avg_tpcp_health, avg_cbs_health",
            "• weakest_dimension",
            "   (base / key / pacing / emerging —",
            "    the block-level rollup with",
            "    the lowest team mean)",
            "",
            "ANSWERS:",
            "• How did my team perform?",
            "• Enough Staff/Principal?",
            "• Who hasn't completed TPCP?",
            "• Which block is weakest?",
        ],
    },
    {
        "title": "B. Individual Drill-Down", "x": Inches(7.55), "color": LIGHT_BLUE, "border": PRIMARY_BLUE,
        "body": [
            "Per-employee flags and lists:",
            "",
            "• low_performers",
            "   combined_competency < 0.50",
            "",
            "• promotable",
            "   combined ≥ 0.75",
            "   AND tenure_readiness ≥ 1.0",
            "   AND passed_tpcp = 'Y'",
            "   AND grade_fit ∈ {0.5, 1.0}",
            "",
            "• close_in_6m",
            "   0.65 ≤ combined < 0.75",
            "",
            "• not_assessed",
            "   has_tpcp = false",
            "",
            "ANSWERS:",
            "• Who is promotable?",
            "• Who are the low performers?",
            "• Who is close in 6 months?",
        ],
    },
]
box_w = Inches(4.8)
box_h = Inches(3.85)
for br in branches:
    x = br["x"]
    y = Inches(3.4)
    add_shape(slide, x, y, box_w, box_h, br["color"], br["border"])
    add_text(slide, x + Inches(0.15), y + Inches(0.15), box_w - Inches(0.3), Inches(0.45),
             br["title"], size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_bullets(slide, x + Inches(0.2), y + Inches(0.65), box_w - Inches(0.4), box_h - Inches(0.75),
                br["body"], size=10.5)


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — Engine Response contract (the JSON shape)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Engine Response — The JSON Contract",
           "One structured document per query — no prose, only enumerated tokens and numbers")

# Left: contract summary
add_shape(slide, Inches(0.4), Inches(1.7), Inches(4.2), Inches(5.3), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(0.6), Inches(1.8), Inches(3.9), Inches(0.4),
         "TOP-LEVEL FIELDS", size=13, bold=True, color=PRIMARY_BLUE)
fields = [
    "• scope",
    "    business, opu, division, team_size",
    "",
    "• join_summary",
    "    both / hris_only / tpcp_only counts",
    "",
    "• team_aggregate",
    "    pass_rate, level_distribution,",
    "    avg_combined, avg_tpcp, avg_cbs,",
    "    weakest_dimension",
    "",
    "• individual_drill_down",
    "    low_performers, promotable,",
    "    close_in_6m, not_assessed",
    "",
    "• employee_health[]",
    "    one 5-score vector per person",
    "",
    "NO gap_matrix field",
    "NO actions field",
]
add_bullets(slide, Inches(0.6), Inches(2.3), Inches(3.9), Inches(4.6), fields, size=11)

# Right: example JSON skeleton
add_shape(slide, Inches(4.9), Inches(1.7), Inches(8.0), Inches(5.3), DARK_BLUE, DARK_BLUE)
add_text(slide, Inches(5.05), Inches(1.8), Inches(7.7), Inches(0.4),
         "EXAMPLE ENGINE_RESPONSE", size=13, bold=True, color=WHITE)
tb = slide.shapes.add_textbox(Inches(5.05), Inches(2.25), Inches(7.7), Inches(4.7))
tf = tb.text_frame
tf.word_wrap = True
example = """{
  "scope": {"business":"Upstream","opu":"PCSB-DR","division":"Drilling","team_size":47},
  "join_summary": {"both":42, "hris_only":4, "tpcp_only":1},
  "team_aggregate": {
    "assessed_count":42, "not_assessed_count":5,
    "pass_rate":0.7619, "level_distribution":{"Staff":18,"Principal":20,"Custodian":4,"Not TP":0},
    "avg_combined_competency":0.68, "avg_tpcp_health":0.72, "avg_cbs_health":0.63,
    "weakest_dimension":"pacing"
  },
  "individual_drill_down": {
    "low_performers":[{"dummy_id":30112,"combined_competency":0.41}],
    "promotable":[{"dummy_id":23267,"combined_competency":0.77}],
    "close_in_6m":[{"dummy_id":28504,"combined_competency":0.68}],
    "not_assessed":[{"dummy_id":40301,"sma_completion_status":"N"}]
  },
  "employee_health":[
    {"dummy_id":23267,"tpcp_health":0.87,"cbs_health":0.65,
     "combined_competency":0.77,"grade_fit":1.0,"tenure_readiness":1.0},
    ...  47 entries total
  ]
}"""
p = tf.paragraphs[0]
p.text = example
p.font.size = Pt(9.5)
p.font.name = "Consolas"
p.font.color.rgb = WHITE


# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — Full request trace (engine stops at diagnosis)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Full Request Trace",
           "One user question walked through every layer — engine returns diagnostic JSON")

steps = [
    ("HOD asks",
     '"Which of my principals in PCSB-DR Drilling\nare promotable, and who is falling behind?"',
     PALE_PURPLE, PURPLE),
    ("Chat parser extracts",
     "business=Upstream, opu=PCSB-DR, division=Drilling,\nquestion_type=promotability+low_performers, filter=assessment_level:Principal",
     PALE_BLUE, PRIMARY_BLUE),
    ("API routes",
     "POST /team/analyze  →  Engine.run(scope)",
     LIGHT_BLUE, PRIMARY_BLUE),
    ("Engine — load + join",
     "HRIS ⋈ TPCP full-outer-join on dummy_id; flag has_hris / has_tpcp",
     PALE_TEAL, TEAL),
    ("Engine — scope filter",
     "Keep rows where business/opu/division match → 47 candidates",
     PALE_RED, RED),
    ("Engine — per-person scoring",
     "Compute 5 health scores for each of 47 employees",
     LIGHT_BLUE, PRIMARY_BLUE),
    ("Engine — analysis",
     "Team aggregate + individual drill-down → 3 promotable, 5 low performers",
     YELLOW_BG, AMBER),
    ("Engine returns JSON; conversation layer narrates",
     '{"individual_drill_down":{...}} → "3 of your principals are promotable; 5 fall below 0.50."',
     PALE_GREEN, GREEN),
]

y = Inches(1.5)
step_h = Inches(0.68)
for i, (t, d, c, b) in enumerate(steps):
    yi = y + i * (step_h + Inches(0.06))
    add_shape(slide, Inches(0.5), yi, Inches(0.7), step_h, b, b, shape=MSO_SHAPE.OVAL)
    add_text(slide, Inches(0.5), yi + Inches(0.18), Inches(0.7), Inches(0.4),
             str(i + 1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_shape(slide, Inches(1.4), yi, Inches(11.4), step_h, c, b)
    add_text(slide, Inches(1.6), yi + Inches(0.06), Inches(4.5), Inches(0.3),
             t, size=12, bold=True, color=DARK_BLUE)
    add_text(slide, Inches(1.6), yi + Inches(0.33), Inches(11), Inches(0.33),
             d, size=10.5, color=DARK_BLUE, font="Consolas")


# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — Out of Scope (explicit: both descopes)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Out of Scope",
           "Two explicit descopes: per-cell competency analysis and action triage")

# Left: per-cell competency descope
add_shape(slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(5.3), PALE_RED, RED)
add_text(slide, Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.4),
         "PER-CELL COMPETENCY — DROPPED", size=14, bold=True, color=RED)
cell_descope = [
    "Engine does NOT consume these TPCP fields:",
    "  • b1..b25, b110 (Base cells)",
    "  • k1..k15, k135 (Key cells)",
    "  • p1..p10, p150 (Pacing cells)",
    "  • e1..e10, e160 (Emerging cells)",
    "",
    "Dropped outputs:",
    "  • gap_matrix (cells / repeat_gaps)",
    "  • critical_gap_score per employee",
    "  • weakest_cells per employee",
    "  • structural / mixed / people labels",
    "",
    "Surviving TPCP signal:",
    "  • base_pct, key_pct, pacing_pct,",
    "    emerging_pct, cti_met_pct",
    "  • team_aggregate.weakest_dimension",
]
add_bullets(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.6), cell_descope, size=11)

# Right: action triage descope
add_shape(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.3), PALE_RED, RED)
add_text(slide, Inches(7.0), Inches(1.8), Inches(5.6), Inches(0.4),
         "ACTION RECOMMENDATIONS — DEFERRED", size=14, bold=True, color=RED)
action_descope = [
    "Engine does NOT emit:",
    "  • actions field in Engine_Response",
    "  • training / exposure / hiring objects",
    "  • priority or horizon_months fields",
    "  • senior-coverage hiring signal",
    "",
    "Questions deferred to a future layer:",
    "  • What actions to close these gaps?",
    "  • Training, exposure, or hiring?",
    "  • Which intervention gives highest",
    "    impact?",
    "  • What's realistic in 6–12 months?",
    "",
    "Downstream triage can still read the",
    "diagnostic JSON (team_aggregate +",
    "individual_drill_down) and apply its",
    "own rules.",
]
add_bullets(slide, Inches(7.0), Inches(2.3), Inches(5.6), Inches(4.6), action_descope, size=11)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Summary — Architecture & Scoring Flow")

summary = [
    "• 5-layer architecture: Presentation → Conversation → API → Engine → Data",
    "• Engine consumes 2 files (HRIS + TPCP), full-outer-joined on dummy_id",
    "• 4-stage engine pipeline: load+join → scope filter → per-person scoring → analysis",
    "• 5 health scores per employee — all deterministic, all in 0–1 range",
    "• 2 parallel analyses: team aggregate + individual drill-down",
    "• TPCP input restricted to block-level rollups (base/key/pacing/emerging %, cti_met %)",
    "• Engine stops at diagnosis — action triage is explicitly out of scope",
    "• Engine emits PURE JSON; conversation layer handles narrative framing",
    "• Every output is traceable back to raw HRIS and TPCP rollup fields",
]
add_bullets(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(5.2), summary, size=16)

add_shape(slide, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.7), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(0.8), Inches(6.6), Inches(11.8), Inches(0.5),
         "Rule-based · Deterministic · Explainable · No labeled data needed · Conversation-ready diagnostic JSON",
         size=13, bold=True, color=PRIMARY_BLUE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
prs.save(str(OUTPUT_FILE))
print(f"Presentation saved to: {OUTPUT_FILE}")
print(f"Total slides: {len(prs.slides)}")
