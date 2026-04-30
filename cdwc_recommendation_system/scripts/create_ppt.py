"""
Generate the CDWC Talent Recommendation Engine presentation deck.
Output: presentation_ppt/CDWC_Talent_Recommendation_Engine.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "presentation_ppt"
OUTPUT_DIR.mkdir(exist_ok=True)
OUTPUT_FILE = OUTPUT_DIR / "CDWC_Talent_Recommendation_Engine.pptx"

# Brand colors
DARK_BLUE = RGBColor(0x1E, 0x29, 0x3B)
PRIMARY_BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
YELLOW_BG = RGBColor(0xFE, 0xF9, 0xC3)
AMBER = RGBColor(0xCA, 0x8A, 0x04)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_BLUE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, left, top, width, height, items, font_size=16, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "CDWC Talent Recommendation Engine",
             font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             "Rule-Based Heuristic Deterministic Computation Method",
             font_size=22, color=RGBColor(0x93, 0xC5, 0xFD), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "Proof of Concept  |  Technical Overview & Product Walkthrough",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Agenda", font_size=32, bold=True, color=DARK_BLUE)

agenda_items = [
    "1.  What is a Rule-Based Heuristic System?",
    "2.  Why Rule-Based over Machine Learning?",
    "3.  System Architecture & Data Flow",
    "4.  The Recommendation Pipeline (4 Stages)",
    "5.  Scoring Functions & Weighted Aggregation",
    "6.  Front-End Product: Form-Based UI & Chat UI",
    "7.  API Endpoints & Integration",
    "8.  Evaluation Framework & POC Sign-Off Criteria",
    "9.  Future Roadmap: Path to Machine Learning",
]
add_bullet_slide_content(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(5),
                         agenda_items, font_size=20, color=DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is a Rule-Based Heuristic System?
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "What is a Rule-Based Heuristic System?", font_size=32, bold=True, color=DARK_BLUE)

# Left column — definition
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8), LIGHT_BLUE, PRIMARY_BLUE)
add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5),
             "Definition", font_size=20, bold=True, color=PRIMARY_BLUE)

defn_items = [
    "• A heuristic is a practical shortcut that produces a 'good enough' answer",
    "• Domain knowledge is encoded as explicit rules and mathematical formulas",
    "• No machine learning, no training, no learned parameters",
    "• Fully deterministic: same input → same output, every time",
]
add_bullet_slide_content(slide, Inches(1.1), Inches(2.5), Inches(5), Inches(3.5),
                         defn_items, font_size=15, color=DARK_BLUE)

# Right column — AI taxonomy
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(4.8), RGBColor(0xF3, 0xE8, 0xFF), RGBColor(0x7C, 0x3A, 0xED))
add_text_box(slide, Inches(7.1), Inches(1.8), Inches(5), Inches(0.5),
             "Where It Fits in AI", font_size=20, bold=True, color=RGBColor(0x7C, 0x3A, 0xED))

taxonomy = [
    "Artificial Intelligence",
    "  ├─ Symbolic AI (rule-based) ← THIS SYSTEM",
    "  │   ├─ Expert Systems",
    "  │   └─ Heuristic Search",
    "  └─ Statistical AI (data-driven)",
    "      ├─ Machine Learning",
    "      └─ Deep Learning / LLMs",
    "",
    "All ML is AI, but not all AI is ML.",
]
add_bullet_slide_content(slide, Inches(7.1), Inches(2.5), Inches(5), Inches(3.5),
                         taxonomy, font_size=14, color=DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Why Rule-Based over Machine Learning?
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Why Rule-Based over Machine Learning?", font_size=32, bold=True, color=DARK_BLUE)

# Comparison table as shapes
headers = ["Factor", "Rule-Based (This System)", "Machine Learning"]
rows = [
    ["Labeled data required?", "No", "Yes (thousands of examples)"],
    ["Training infrastructure?", "None", "Training pipeline + compute"],
    ["Explainability", "Full — every score traceable", "Varies (black box for neural nets)"],
    ["Time to build", "Fast — encode domain knowledge", "Slower — data collection + tuning"],
    ["Maintenance", "Update rules directly", "Retrain, monitor for drift"],
    ["Deterministic?", "Yes — always reproducible", "Depends on implementation"],
]

col_widths = [Inches(3), Inches(4), Inches(4)]
row_h = Inches(0.55)
start_x = Inches(0.8)
start_y = Inches(1.8)

# Header row
x = start_x
for i, h in enumerate(headers):
    s = add_shape(slide, x, start_y, col_widths[i], row_h, PRIMARY_BLUE)
    s.text_frame.paragraphs[0].text = h
    s.text_frame.paragraphs[0].font.size = Pt(14)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    s.text_frame.paragraphs[0].font.name = "Calibri"
    x += col_widths[i] + Inches(0.05)

# Data rows
for ri, row in enumerate(rows):
    y = start_y + (ri + 1) * (row_h + Inches(0.05))
    x = start_x
    bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
    for ci, cell in enumerate(row):
        s = add_shape(slide, x, y, col_widths[ci], row_h, bg, RGBColor(0xE2, 0xE8, 0xF0))
        s.text_frame.paragraphs[0].text = cell
        s.text_frame.paragraphs[0].font.size = Pt(13)
        s.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        s.text_frame.paragraphs[0].font.name = "Calibri"
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        x += col_widths[ci] + Inches(0.05)

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.6),
             "Key insight: When labeled data is unavailable, rule-based is the recommended approach. "
             "It serves as both a working system and a baseline for future ML.",
             font_size=15, bold=True, color=PRIMARY_BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — System Architecture & Data Flow
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "System Architecture & Data Flow", font_size=32, bold=True, color=DARK_BLUE)

# Architecture boxes
components = [
    {"label": "Frontend\n(Vanilla JS / Streamlit)", "x": 0.5, "y": 2.0, "w": 2.5, "h": 1.5, "color": RGBColor(0xF3, 0xE8, 0xFF), "border": RGBColor(0x7C, 0x3A, 0xED)},
    {"label": "FastAPI\nBackend", "x": 3.8, "y": 2.0, "w": 2.2, "h": 1.5, "color": RGBColor(0xDB, 0xEA, 0xFE), "border": PRIMARY_BLUE},
    {"label": "Recommendation\nEngine", "x": 6.8, "y": 2.0, "w": 2.5, "h": 1.5, "color": RGBColor(0xFE, 0xF9, 0xC3), "border": AMBER},
    {"label": "Employee\nData (JSON)", "x": 10.1, "y": 2.0, "w": 2.2, "h": 1.5, "color": RGBColor(0xD1, 0xFA, 0xE5), "border": GREEN},
]

for c in components:
    s = add_shape(slide, Inches(c["x"]), Inches(c["y"]), Inches(c["w"]), Inches(c["h"]), c["color"], c["border"])
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = c["label"]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    tf.paragraphs[0].space_before = Pt(15)

# Arrow labels
arrows = [
    {"text": "HTTP →", "x": 3.1, "y": 2.5},
    {"text": "Function call →", "x": 6.1, "y": 2.5},
    {"text": "Load →", "x": 9.4, "y": 2.5},
]
for a in arrows:
    add_text_box(slide, Inches(a["x"]), Inches(a["y"]), Inches(1.5), Inches(0.4),
                 a["text"], font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Data sources section
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
             "Data Sources Consumed:", font_size=18, bold=True, color=DARK_BLUE)

data_sources = [
    "• EPDR (Employee Performance & Development Review) — competency assessments, proficiency levels",
    "• SMA HRIS — skill groups, job/salary grades, CBS percentages, completion status",
    "• Item Catalog Report — competency inventories (technical & behavioral)",
    "• HR Flex Report — supplementary HR data",
]
add_bullet_slide_content(slide, Inches(1.0), Inches(4.8), Inches(11), Inches(2.5),
                         data_sources, font_size=15, color=DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — The Recommendation Pipeline (4 Stages)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "The Recommendation Pipeline", font_size=32, bold=True, color=DARK_BLUE)

stages = [
    {
        "title": "Stage 1: Hard Filtering",
        "desc": "Remove candidates who fail\nmandatory constraints:\n• Unavailable → rejected\n• Role level too low → rejected",
        "color": RGBColor(0xFE, 0xE2, 0xE2), "border": RED,
    },
    {
        "title": "Stage 2: Feature Scoring",
        "desc": "Compute 0–1 similarity scores:\n• Skill overlap (Jaccard)\n• Competency similarity\n• Experience similarity\n• Role match score",
        "color": LIGHT_BLUE, "border": PRIMARY_BLUE,
    },
    {
        "title": "Stage 3: Weighted\nAggregation",
        "desc": "Combine scores:\n  skill_overlap × 0.35\n+ competency × 0.25\n+ experience × 0.20\n+ role_match × 0.20\n= total_score",
        "color": YELLOW_BG, "border": AMBER,
    },
    {
        "title": "Stage 4: Ranking",
        "desc": "Sort candidates by\ntotal_score descending.\n\nReturn Top-K results\n(default K=5) with full\nscore breakdowns.",
        "color": RGBColor(0xD1, 0xFA, 0xE5), "border": GREEN,
    },
]

box_w = Inches(2.7)
box_h = Inches(4.2)
gap = Inches(0.3)
start_x = Inches(0.6)

for i, st in enumerate(stages):
    x = start_x + i * (box_w + gap)
    y = Inches(1.6)

    s = add_shape(slide, x, y, box_w, box_h, st["color"], st["border"])

    # Title
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), box_w - Inches(0.4), Inches(0.8),
                 st["title"], font_size=16, bold=True, color=DARK_BLUE)

    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), box_w - Inches(0.4), Inches(2.8),
                 st["desc"], font_size=13, color=DARK_BLUE)

# Arrow connectors between stages
for i in range(3):
    x = start_x + (i + 1) * (box_w + gap) - gap + Inches(0.02)
    add_text_box(slide, x, Inches(3.4), gap, Inches(0.4),
                 "→", font_size=24, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Scoring Functions Deep Dive
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Scoring Functions & Weighted Aggregation", font_size=32, bold=True, color=DARK_BLUE)

# Table
headers = ["Function", "Formula", "What It Measures"]
rows = [
    ["Skill Overlap", "|intersection| / |required|", "How many required skills the candidate has"],
    ["Competency\nSimilarity", "1 − (|candidate − required| / 5.0)", "How close competency is to requirement"],
    ["Experience\nSimilarity", "min(candidate_yrs / required, 1.0)", "Whether experience meets the minimum"],
    ["Role Match", "1.0 exact | 0.5 overqualified | 0.0 under", "Seniority alignment"],
]

col_widths = [Inches(2.5), Inches(4.5), Inches(4.5)]
row_h = Inches(0.75)
start_x = Inches(0.6)
start_y = Inches(1.6)

x = start_x
for i, h in enumerate(headers):
    s = add_shape(slide, x, start_y, col_widths[i], Inches(0.55), PRIMARY_BLUE)
    s.text_frame.paragraphs[0].text = h
    s.text_frame.paragraphs[0].font.size = Pt(14)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = WHITE
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    s.text_frame.paragraphs[0].font.name = "Calibri"
    x += col_widths[i] + Inches(0.05)

for ri, row in enumerate(rows):
    y = start_y + Inches(0.6) + ri * (row_h + Inches(0.05))
    x = start_x
    bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
    for ci, cell in enumerate(row):
        s = add_shape(slide, x, y, col_widths[ci], row_h, bg, RGBColor(0xE2, 0xE8, 0xF0))
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cell
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_BLUE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        x += col_widths[ci] + Inches(0.05)

# Weight formula box
add_shape(slide, Inches(0.6), Inches(5.2), Inches(11.5), Inches(1.5), RGBColor(0xF0, 0xFD, 0xF4), GREEN)
add_text_box(slide, Inches(0.9), Inches(5.3), Inches(11), Inches(0.4),
             "Weighted Aggregation Formula:", font_size=16, bold=True, color=GREEN)
add_text_box(slide, Inches(0.9), Inches(5.8), Inches(11), Inches(0.7),
             "Total Score = (0.35 × Skill Overlap) + (0.25 × Competency) + (0.20 × Experience) + (0.20 × Role Match)",
             font_size=18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Front-End Product: Form-Based UI
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Front-End Product: Form-Based Search UI", font_size=32, bold=True, color=DARK_BLUE)

# Left — UI mockup
add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2), RGBColor(0xF8, 0xFA, 0xFC), RGBColor(0xE2, 0xE8, 0xF0))
add_text_box(slide, Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.4),
             "Form Search UI — localhost:8000/ui", font_size=14, bold=True, color=PRIMARY_BLUE)

mockup_lines = [
    "┌─────────────────────────────────────────┐",
    "│  Skills: [python, machine_learning, sql] │",
    "│                                          │",
    "│  Min Competency:  [====●=====] 3.0       │",
    "│  Min Experience:  [==●=======] 2 yrs     │",
    "│                                          │",
    "│  Role Level:  ○ Junior  ● Mid  ○ Senior  │",
    "│  ☑ Available only                        │",
    "│                                          │",
    "│  [ 🔍 Search Talent ]                    │",
    "└─────────────────────────────────────────┘",
]
add_bullet_slide_content(slide, Inches(1.0), Inches(2.3), Inches(5.2), Inches(4),
                         mockup_lines, font_size=12, color=DARK_BLUE)

# Right — features
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), LIGHT_BLUE, PRIMARY_BLUE)
add_text_box(slide, Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4),
             "Key Features", font_size=20, bold=True, color=PRIMARY_BLUE)

features = [
    "• Pure vanilla HTML/CSS/JS — zero framework dependencies",
    "• Served directly by FastAPI at /ui endpoint",
    "• Two search modes: Form-based & Chat-based",
    "• Form mode calls POST /recommend directly",
    "• Chat mode calls POST /chat (NLP parsing)",
    "• Real-time result cards with score breakdowns",
    "• Skill tags, department info, experience display",
    "• Responsive design — works on mobile & desktop",
    "• CORS-enabled for cross-origin deployment",
    "• Role levels loaded dynamically from API",
]
add_bullet_slide_content(slide, Inches(7.1), Inches(2.5), Inches(5.2), Inches(4),
                         features, font_size=14, color=DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Front-End Product: Result Cards
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Front-End Product: Result Cards & Chat Mode", font_size=32, bold=True, color=DARK_BLUE)

# Result card mockup
add_shape(slide, Inches(0.6), Inches(1.6), Inches(6), Inches(4.8), RGBColor(0xF8, 0xFA, 0xFC), RGBColor(0xE2, 0xE8, 0xF0))
add_text_box(slide, Inches(0.9), Inches(1.7), Inches(5.5), Inches(0.4),
             "Result Cards", font_size=16, bold=True, color=PRIMARY_BLUE)

# Card 1
add_shape(slide, Inches(1.0), Inches(2.3), Inches(5.2), Inches(1.6), WHITE, PRIMARY_BLUE)
card1 = [
    "1. Alice Johnson — Score: 0.87",
    "   Engineering · Senior · 8 yrs experience",
    "   Skills: python | ml | sql | tensorflow",
    "   Skill 92% | Competency 85% | Exp 80% | Role 100%",
]
add_bullet_slide_content(slide, Inches(1.2), Inches(2.4), Inches(4.8), Inches(1.4),
                         card1, font_size=12, color=DARK_BLUE)

# Card 2
add_shape(slide, Inches(1.0), Inches(4.1), Inches(5.2), Inches(1.6), WHITE, PRIMARY_BLUE)
card2 = [
    "2. Bob Smith — Score: 0.74",
    "   Data Science · Mid · 5 yrs experience",
    "   Skills: python | sql | pandas",
    "   Skill 78% | Competency 70% | Exp 72% | Role 50%",
]
add_bullet_slide_content(slide, Inches(1.2), Inches(4.2), Inches(4.8), Inches(1.4),
                         card2, font_size=12, color=DARK_BLUE)

# Chat mode description
add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(4.8), RGBColor(0xF3, 0xE8, 0xFF), RGBColor(0x7C, 0x3A, 0xED))
add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5.2), Inches(0.4),
             "Chat Search Mode", font_size=16, bold=True, color=RGBColor(0x7C, 0x3A, 0xED))

chat_desc = [
    "Natural language input — same as Streamlit:",
    "",
    '  User: "I need a senior Python developer',
    '         with ML experience, 5+ years"',
    "",
    "  System parses → extracts structured query →",
    "  calls recommendation engine → returns results",
    "",
    "Two parser backends available:",
    "  • MockParser — keyword-based, zero dependencies",
    "  • LLMParser — OpenAI-powered (optional)",
]
add_bullet_slide_content(slide, Inches(7.3), Inches(2.5), Inches(5.2), Inches(3.5),
                         chat_desc, font_size=13, color=DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — API Endpoints
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "API Endpoints & Integration", font_size=32, bold=True, color=DARK_BLUE)

endpoints = [
    ["POST /recommend", "Core endpoint — accepts structured project requirements, returns ranked candidates"],
    ["POST /chat", "Natural language endpoint — accepts plain English, parses and returns recommendations"],
    ["GET /health", "Liveness check for monitoring and deployment"],
    ["GET /history", "Paginated audit trail of all recommendation requests"],
    ["GET /history/{id}", "Retrieve a specific historical recommendation by request ID"],
    ["GET /config/weights", "View current scoring weights"],
    ["PUT /config/weights", "Update scoring weights at runtime (must sum to 1.0)"],
    ["GET /config/role-levels", "List available role levels for UI dropdowns"],
    ["GET /docs", "Auto-generated interactive API documentation (Swagger UI)"],
]

row_h = Inches(0.55)
start_y = Inches(1.5)
col1_w = Inches(3.2)
col2_w = Inches(8.5)

# Header
s = add_shape(slide, Inches(0.6), start_y, col1_w, row_h, PRIMARY_BLUE)
s.text_frame.paragraphs[0].text = "Endpoint"
s.text_frame.paragraphs[0].font.size = Pt(14)
s.text_frame.paragraphs[0].font.bold = True
s.text_frame.paragraphs[0].font.color.rgb = WHITE
s.text_frame.paragraphs[0].font.name = "Calibri"
s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

s = add_shape(slide, Inches(0.6) + col1_w + Inches(0.05), start_y, col2_w, row_h, PRIMARY_BLUE)
s.text_frame.paragraphs[0].text = "Description"
s.text_frame.paragraphs[0].font.size = Pt(14)
s.text_frame.paragraphs[0].font.bold = True
s.text_frame.paragraphs[0].font.color.rgb = WHITE
s.text_frame.paragraphs[0].font.name = "Calibri"
s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

for ri, (ep, desc) in enumerate(endpoints):
    y = start_y + (ri + 1) * (row_h + Inches(0.03))
    bg = LIGHT_BLUE if ri % 2 == 0 else WHITE

    s = add_shape(slide, Inches(0.6), y, col1_w, row_h, bg, RGBColor(0xE2, 0xE8, 0xF0))
    s.text_frame.paragraphs[0].text = ep
    s.text_frame.paragraphs[0].font.size = Pt(12)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    s.text_frame.paragraphs[0].font.name = "Consolas"
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    s = add_shape(slide, Inches(0.6) + col1_w + Inches(0.05), y, col2_w, row_h, bg, RGBColor(0xE2, 0xE8, 0xF0))
    s.text_frame.paragraphs[0].text = desc
    s.text_frame.paragraphs[0].font.size = Pt(12)
    s.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    s.text_frame.paragraphs[0].font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Evaluation Framework
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Evaluation Framework & Metrics", font_size=32, bold=True, color=DARK_BLUE)

# Metrics table
metrics = [
    ["Precision@K", "Of the top K recommendations, how many are relevant?"],
    ["Recall@K", "Of all relevant candidates, how many appear in top K?"],
    ["NDCG@K", "Are relevant candidates ranked higher? (position-aware)"],
    ["Diversity", "Are recommendations from varied departments?"],
    ["Fairness", "Does department representation reflect the eligible pool?"],
    ["Ranking Stability", "Do rankings hold under small weight perturbations? (Kendall's τ)"],
]

start_y = Inches(1.5)
for ri, (metric, desc) in enumerate(metrics):
    y = start_y + ri * Inches(0.7)
    bg = LIGHT_BLUE if ri % 2 == 0 else WHITE

    s = add_shape(slide, Inches(0.6), y, Inches(3), Inches(0.6), bg, RGBColor(0xE2, 0xE8, 0xF0))
    s.text_frame.paragraphs[0].text = metric
    s.text_frame.paragraphs[0].font.size = Pt(14)
    s.text_frame.paragraphs[0].font.bold = True
    s.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    s.text_frame.paragraphs[0].font.name = "Calibri"
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    s = add_shape(slide, Inches(3.65), y, Inches(8.8), Inches(0.6), bg, RGBColor(0xE2, 0xE8, 0xF0))
    s.text_frame.paragraphs[0].text = desc
    s.text_frame.paragraphs[0].font.size = Pt(13)
    s.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    s.text_frame.paragraphs[0].font.name = "Calibri"

# Stability note
add_shape(slide, Inches(0.6), Inches(5.8), Inches(11.8), Inches(1.2), RGBColor(0xFE, 0xF9, 0xC3), AMBER)
add_text_box(slide, Inches(0.9), Inches(5.9), Inches(11.2), Inches(0.4),
             "Stability Testing", font_size=16, bold=True, color=AMBER)
add_text_box(slide, Inches(0.9), Inches(6.3), Inches(11.2), Inches(0.5),
             "Weights are perturbed by ±5% across 10 iterations. Kendall's Tau measures ranking consistency. "
             "High stability (τ ≈ 1.0) confirms the system is robust and not overly sensitive to exact weight values.",
             font_size=13, color=DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — POC Sign-Off Criteria
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "POC Sign-Off Criteria", font_size=32, bold=True, color=DARK_BLUE)

# Three columns
categories = [
    {
        "title": "Quantitative",
        "items": [
            "☐ Precision@5 ≥ 70% on expert-\n   reviewed scenarios (min 25)",
            "☐ Hit Rate@5 ≥ 60% on historical\n   assignment data (if available)",
            "☐ Outperforms random baseline\n   by ≥ 2x on Precision@5",
            "☐ Ranking Stability (τ) ≥ 0.8",
        ],
        "color": LIGHT_BLUE, "border": PRIMARY_BLUE,
    },
    {
        "title": "Qualitative",
        "items": [
            "☐ Domain experts agree rankings\n   'make sense' for ≥ 80% of cases",
            "☐ Passes all edge case scenarios\n   without nonsensical results",
            "☐ Response time < 2 seconds\n   for recommendation generation",
        ],
        "color": RGBColor(0xF3, 0xE8, 0xFF), "border": RGBColor(0x7C, 0x3A, 0xED),
    },
    {
        "title": "Operational",
        "items": [
            "☐ Handles full employee dataset\n   without errors",
            "☐ Results are fully reproducible\n   (deterministic)",
            "☐ API documentation complete\n   and accessible at /docs",
        ],
        "color": RGBColor(0xD1, 0xFA, 0xE5), "border": GREEN,
    },
]

col_w = Inches(3.8)
for i, cat in enumerate(categories):
    x = Inches(0.5) + i * (col_w + Inches(0.25))
    y = Inches(1.5)
    add_shape(slide, x, y, col_w, Inches(5), cat["color"], cat["border"])
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.4), Inches(0.5),
                 cat["title"], font_size=18, bold=True, color=DARK_BLUE)
    add_bullet_slide_content(slide, x + Inches(0.2), y + Inches(0.8), col_w - Inches(0.4), Inches(3.8),
                             cat["items"], font_size=13, color=DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — Future Roadmap: Path to ML
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Future Roadmap: Path to Machine Learning", font_size=32, bold=True, color=DARK_BLUE)

# Timeline phases
phases = [
    {
        "title": "Phase 1 (Current)\nRule-Based POC",
        "desc": "• Hand-crafted heuristic scoring\n• Deterministic, explainable\n• No labeled data needed\n• Serves as ML baseline",
        "color": RGBColor(0xD1, 0xFA, 0xE5), "border": GREEN,
    },
    {
        "title": "Phase 2\nData Collection",
        "desc": "• Collect match outcomes\n• Manager ratings on fit\n• Expert panel scoring\n• Build labeled dataset",
        "color": YELLOW_BG, "border": AMBER,
    },
    {
        "title": "Phase 3\nHybrid ML",
        "desc": "• Keep rule-based features\n• ML learns optimal weights\n• XGBoost / LightGBM\n• A/B test vs rule-based",
        "color": LIGHT_BLUE, "border": PRIMARY_BLUE,
    },
    {
        "title": "Phase 4\nFull ML Pipeline",
        "desc": "• End-to-end learned model\n• Learning to Rank\n• Continuous retraining\n• Drift monitoring",
        "color": RGBColor(0xF3, 0xE8, 0xFF), "border": RGBColor(0x7C, 0x3A, 0xED),
    },
]

box_w = Inches(2.7)
box_h = Inches(4)
for i, ph in enumerate(phases):
    x = Inches(0.5) + i * (box_w + Inches(0.3))
    y = Inches(1.8)
    add_shape(slide, x, y, box_w, box_h, ph["color"], ph["border"])
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), box_w - Inches(0.4), Inches(0.9),
                 ph["title"], font_size=15, bold=True, color=DARK_BLUE)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.3), box_w - Inches(0.4), Inches(2.5),
                 ph["desc"], font_size=13, color=DARK_BLUE)

# Arrows
for i in range(3):
    x = Inches(0.5) + (i + 1) * (box_w + Inches(0.3)) - Inches(0.28)
    add_text_box(slide, x, Inches(3.5), Inches(0.3), Inches(0.4),
                 "→", font_size=24, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
             "Recommended hybrid approach: keep existing features as engineered inputs, use ML to learn optimal weights — "
             "replacing the hand-tuned 0.35/0.25/0.20/0.20 with data-driven values.",
             font_size=14, bold=True, color=PRIMARY_BLUE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You / Q&A
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
             "Thank You", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             "Questions & Discussion", font_size=24, color=RGBColor(0x93, 0xC5, 0xFD), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
             "API Docs: http://localhost:8000/docs   |   UI: http://localhost:8000/ui",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
prs.save(str(OUTPUT_FILE))
print(f"Presentation saved to: {OUTPUT_FILE}")
print(f"Total slides: {len(prs.slides)}")
