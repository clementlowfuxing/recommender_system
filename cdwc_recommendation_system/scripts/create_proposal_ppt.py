"""
Generate the CDWC Talent Diagnostic Engine — Redesign Proposal deck.
Output: presentation_ppt/CDWC_Redesign_Proposal.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "presentation_ppt"
OUTPUT_DIR.mkdir(exist_ok=True)
OUTPUT_FILE = OUTPUT_DIR / "CDWC_Redesign_Proposal.pptx"

# Palette
DARK_BLUE = RGBColor(0x1E, 0x29, 0x3B)
PRIMARY_BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
PALE_BLUE = RGBColor(0xF0, 0xF7, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
GREEN = RGBColor(0x05, 0x96, 0x69)
PALE_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
RED = RGBColor(0xDC, 0x26, 0x26)
PALE_RED = RGBColor(0xFE, 0xE2, 0xE2)
YELLOW_BG = RGBColor(0xFE, 0xF9, 0xC3)
AMBER = RGBColor(0xCA, 0x8A, 0x04)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PALE_PURPLE = RGBColor(0xF3, 0xE8, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK_BLUE, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def add_bullets(slide, left, top, width, height, items, size=15, color=DARK_BLUE, bold_first=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        if bold_first and i == 0:
            p.font.bold = True
    return tb


def add_header(slide, title, subtitle=None):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7),
             title, size=28, bold=True, color=DARK_BLUE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GRAY)
    # Accent bar
    add_shape(slide, Inches(0.6), Inches(1.35), Inches(0.7), Inches(0.05), PRIMARY_BLUE)


def add_table(slide, left, top, col_widths, row_h, headers, rows, header_color=PRIMARY_BLUE,
              alt_bg=LIGHT_BLUE, header_font=13, body_font=12):
    x0 = left
    y = top
    # headers
    x = x0
    for i, h in enumerate(headers):
        s = add_shape(slide, x, y, col_widths[i], Inches(0.5), header_color)
        s.text_frame.paragraphs[0].text = h
        s.text_frame.paragraphs[0].font.size = Pt(header_font)
        s.text_frame.paragraphs[0].font.bold = True
        s.text_frame.paragraphs[0].font.color.rgb = WHITE
        s.text_frame.paragraphs[0].font.name = "Calibri"
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        x += col_widths[i] + Inches(0.04)
    # rows
    for ri, row in enumerate(rows):
        y_r = y + Inches(0.55) + ri * (row_h + Inches(0.04))
        x = x0
        bg = alt_bg if ri % 2 == 0 else WHITE
        for ci, cell in enumerate(row):
            s = add_shape(slide, x, y_r, col_widths[ci], row_h, bg, RGBColor(0xE2, 0xE8, 0xF0))
            tf = s.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(cell)
            p.font.size = Pt(body_font)
            p.font.color.rgb = DARK_BLUE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.LEFT
            x += col_widths[ci] + Inches(0.04)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text(slide, Inches(1), Inches(1.6), Inches(11), Inches(1.2),
         "CDWC Talent Diagnostic Engine",
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
         "Redesign Proposal",
         size=26, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
         "From Project-Matching to Team Diagnostics, Gap Analysis & Action Insights",
         size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
         "Data-Driven Redesign Based on TPCP + SMA HRIS",
         size=14, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Agenda")

items = [
    "1.  The Shift — From Matching to Diagnosing",
    "2.  What Users Actually Ask (16 HOD/GM Questions)",
    "3.  What the Data Gives Us (TPCP + SMA HRIS)",
    "4.  Why the Current Scoring Functions No Longer Fit",
    "5.  New Core Primitives — 6 Health Scores",
    "6.  Hard Filters & Org Scope",
    "7.  New Pipeline Architecture",
    "8.  Output for Each Question Type",
    "9.  Gap Analysis — People vs Structural",
    "10. Action Recommendation Logic",
    "11. New API Surface",
    "12. Module Structure & Implementation Order",
    "13. Migration Impact",
    "14. Summary & Next Steps",
]
add_bullets(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(5.5), items, size=17)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Shift
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "The Shift — From Matching to Diagnosing")

# OLD
add_shape(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(5), PALE_RED, RED)
add_text(slide, Inches(0.9), Inches(1.85), Inches(5.4), Inches(0.5),
         "OLD: Project Matching Engine", size=18, bold=True, color=RED)
old_items = [
    "• Input: project requirements",
    "• Output: ranked candidates",
    "• Core question:",
    "   'Who best fits this project?'",
    "",
    "• Depends on: skills list, experience years,",
    "   role level, availability",
    "• These fields DO NOT EXIST in the",
    "   real HRIS / TPCP data",
    "",
    "• Does not answer any of the actual",
    "   questions HODs and GMs ask",
]
add_bullets(slide, Inches(0.9), Inches(2.4), Inches(5.4), Inches(4.2), old_items, size=14, color=DARK_BLUE)

# NEW
add_shape(slide, Inches(6.9), Inches(1.7), Inches(6), Inches(5), PALE_GREEN, GREEN)
add_text(slide, Inches(7.2), Inches(1.85), Inches(5.4), Inches(0.5),
         "NEW: Talent Diagnostic Engine", size=18, bold=True, color=GREEN)
new_items = [
    "• Input: team scope (business / opu / division)",
    "• Output: team health + individuals + gaps",
    "   + action recommendations (as JSON)",
    "• Core question:",
    "   'How healthy is my team, where are",
    "   the gaps, who is ready, what next?'",
    "",
    "• Uses real HRIS + TPCP fields as-is",
    "• Answers all 16 HOD/GM questions",
    "• Conversation layer converts JSON",
    "   into narrative insights",
]
add_bullets(slide, Inches(7.2), Inches(2.4), Inches(5.4), Inches(4.2), new_items, size=14, color=DARK_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — What users actually ask (grouped into 4 workflows)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "What Users Actually Ask",
           "16 HOD/GM questions map to 4 analytical workflows")

cats = [
    {"title": "A. Individual Drill-Down", "color": LIGHT_BLUE, "border": PRIMARY_BLUE,
     "items": [
         "Who is promotable in my division?",
         "What gaps need to be closed",
         "   before promotion?",
         "Can this talent be moved into",
         "   a critical role?",
     ]},
    {"title": "B. Team Aggregate", "color": PALE_PURPLE, "border": PURPLE,
     "items": [
         "How did my talents perform",
         "   in the latest TPCP?",
         "Who has not completed TPCP?",
         "What is the overall SMA strength",
         "   for my team?",
         "Do we have enough people at",
         "   Staff and Principal level?",
     ]},
    {"title": "C. Gap Analysis", "color": YELLOW_BG, "border": AMBER,
     "items": [
         "Which staff / principal still have",
         "   critical gaps?",
         "Any high-risk critical roles",
         "   from TPCP results?",
         "Which competencies are",
         "   consistently weak?",
         "Is this a people issue or a",
         "   structural competency gap?",
         "Are there repeat gaps across",
         "   multiple talents?",
     ]},
    {"title": "D. Action Recommendations", "color": PALE_GREEN, "border": GREEN,
     "items": [
         "What actions to close these gaps?",
         "Training, exposure, or hiring?",
         "Which capability intervention",
         "   gives the highest impact?",
         "What can be realistically closed",
         "   in the next 6–12 months?",
     ]},
]

w = Inches(3.0)
gap = Inches(0.15)
start_x = Inches(0.6)
for i, c in enumerate(cats):
    x = start_x + i * (w + gap)
    y = Inches(1.75)
    add_shape(slide, x, y, w, Inches(5.4), c["color"], c["border"])
    add_text(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.5),
             c["title"], size=15, bold=True, color=DARK_BLUE)
    add_bullets(slide, x + Inches(0.2), y + Inches(0.75), w - Inches(0.4), Inches(4.6),
                c["items"], size=11)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — What the data gives us
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "What the Data Gives Us",
           "HRIS and TPCP joined on dummy_id")

# Two data source boxes
add_shape(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(3.2), LIGHT_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(0.9), Inches(1.85), Inches(5.4), Inches(0.5),
         "SMA HRIS  —  9,455 rows", size=16, bold=True, color=PRIMARY_BLUE)
hris_fields = [
    "Identity: dummy_id, employee_name, superior_name",
    "Org: skill_group",
    "Grade: job_grade, salary_grade, year_in_sg",
    "Completion: sma_completion_status",
    "CBS Percentages (0–1):",
    "   • overall_cbs_pct",
    "   • technical_cbs_pct",
    "   • leadership_cbs_pct",
]
add_bullets(slide, Inches(0.9), Inches(2.35), Inches(5.4), Inches(2.5), hris_fields, size=12)

add_shape(slide, Inches(6.9), Inches(1.7), Inches(6), Inches(3.2), PALE_PURPLE, PURPLE)
add_text(slide, Inches(7.2), Inches(1.85), Inches(5.4), Inches(0.5),
         "TPCP Assessment  —  1,018 rows", size=16, bold=True, color=PURPLE)
tpcp_fields = [
    "Org: business, opu, division, skill_group,",
    "   discipline, sub_discipline",
    "Level: assessment_level, assessment_type,",
    "   qualified_level, passed_tpcp",
    "Rolled-up %s: base_pct, key_pct, pacing_pct,",
    "   emerging_pct, cti_met_pct",
    "Competency cells: 110+ fields",
    "   (B1–B25, K1–K15, P1–P10, E1–E10 + 2nd block)",
]
add_bullets(slide, Inches(7.2), Inches(2.35), Inches(5.4), Inches(2.5), tpcp_fields, size=12)

# Join summary
add_shape(slide, Inches(0.6), Inches(5.1), Inches(12.3), Inches(1.8), PALE_GREEN, GREEN)
add_text(slide, Inches(0.9), Inches(5.25), Inches(11.7), Inches(0.5),
         "Left-join on dummy_id  →  unified employee records", size=16, bold=True, color=GREEN)
overlap = [
    "• 977 employees have BOTH HRIS and TPCP records  (full signal)",
    "• 8,477 employees are HRIS-only  (scored on CBS alone, zero on TPCP dimensions per decision)",
    "• 41 employees are TPCP-only  (edge case; scored on TPCP alone)",
]
add_bullets(slide, Inches(0.9), Inches(5.75), Inches(11.7), Inches(1.1), overlap, size=13, color=DARK_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Why current scoring functions no longer fit
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Why Current Scoring Functions No Longer Fit")

headers = ["Old Function", "Status", "Why It Fails"]
rows = [
    ["skill_overlap_score", "DELETED",
     "No 'skills' list in either file. Closest field is skill_group (broad category), not a skill-matching problem."],
    ["experience_similarity", "REPLACED",
     "No 'years_experience' field. year_in_sg is re-purposed for promotion readiness, not similarity matching."],
    ["role_match_score", "REPLACED",
     "junior/mid/senior scale does not exist. Real data uses job_grade (P1–P12) and assessment_level (Staff/Principal/Custodian)."],
    ["competency_similarity", "REPLACED",
     "No single 0–5 score. Instead: 3 CBS %s + 4 TPCP dimension %s + pass/fail. Concept shifts from similarity to absolute health."],
]
add_table(slide, Inches(0.6), Inches(1.75),
          [Inches(3.0), Inches(1.6), Inches(7.7)],
          Inches(0.9), headers, rows)

add_shape(slide, Inches(0.6), Inches(6.2), Inches(12.3), Inches(0.9), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(0.9), Inches(6.35), Inches(11.7), Inches(0.6),
         "Framework stays (filter → score → aggregate → insight). "
         "Functions are rewritten against fields that actually exist.",
         size=13, bold=True, color=PRIMARY_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — New core primitives: 6 health scores
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "New Core Primitives — 6 Health Scores (all 0–1)")

headers = ["Score", "Formula (simplified)", "Source Fields"]
rows = [
    ["1. TPCP Health",
     "0.6 × cti_met_pct + 0.4 × avg(base, key, pacing, emerging)\n+ pass boost (capped at 1.0)",
     "TPCP: cti_met_pct, base_pct, key_pct,\npacing_pct, emerging_pct, passed_tpcp"],
    ["2. CBS Health",
     "0.5 × technical_cbs + 0.3 × overall_cbs + 0.2 × leadership_cbs\n(weights rebalanced if a dimension is missing)",
     "HRIS: technical_cbs_pct, overall_cbs_pct,\nleadership_cbs_pct"],
    ["3. Combined Competency",
     "0.55 × TPCP Health + 0.45 × CBS Health\n(fallback to CBS alone if no TPCP record)",
     "Derived from #1 and #2"],
    ["4. Grade Fit",
     "Alignment of job_grade ↔ assessment_level\n(exact=1.0, over=0.5, under=0.0)",
     "HRIS: job_grade  |  TPCP: assessment_level"],
    ["5. Tenure Readiness",
     "min(year_in_sg / grade_tenure_threshold, 1.0)\nthresholds: P1–P3 → 3y, P4–P6 → 4y, P7+ → 5y",
     "HRIS: year_in_sg, job_grade"],
    ["6. Critical Gap Score",
     "1 − (count of competency cells scored 0 or 1 / total cells)\nfrom B/K/P/E blocks",
     "TPCP: B1–B25, K1–K15, P1–P10, E1–E10"],
]
add_table(slide, Inches(0.4), Inches(1.7),
          [Inches(2.3), Inches(5.3), Inches(4.7)],
          Inches(0.8), headers, rows, header_font=12, body_font=11)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Hard Filters & Org Scope
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Hard Filters & Org Scope",
           "What narrows the population before scoring")

# Filter (only one)
add_shape(slide, Inches(0.6), Inches(1.75), Inches(6), Inches(5.3), PALE_RED, RED)
add_text(slide, Inches(0.9), Inches(1.9), Inches(5.4), Inches(0.5),
         "Mandatory Hard Filter", size=18, bold=True, color=RED)
add_shape(slide, Inches(0.9), Inches(2.6), Inches(5.4), Inches(1.5), WHITE, RED)
add_text(slide, Inches(1.1), Inches(2.7), Inches(5), Inches(0.5),
         "Org Scope (match ALL):", size=14, bold=True, color=DARK_BLUE)
add_bullets(slide, Inches(1.2), Inches(3.1), Inches(5), Inches(0.9),
            ["business   =   user.business",
             "opu        =   user.opu",
             "division   =   user.division"], size=13)

# Soft dropped
add_text(slide, Inches(0.9), Inches(4.3), Inches(5.4), Inches(0.5),
         "Dropped / Demoted", size=15, bold=True, color=RED)
dropped = [
    "• availability  →  field does not exist",
    "• skills overlap  →  no skills list in data",
    "• role_level (junior/mid/senior)  →  no such scale",
    "• min_experience  →  no years_experience field",
]
add_bullets(slide, Inches(1.1), Inches(4.9), Inches(5), Inches(2), dropped, size=12, color=DARK_BLUE)

# What becomes signal
add_shape(slide, Inches(6.9), Inches(1.75), Inches(6), Inches(5.3), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(7.2), Inches(1.9), Inches(5.4), Inches(0.5),
         "Everything Else → Signal, Not Filter", size=18, bold=True, color=PRIMARY_BLUE)
signals = [
    "• SMA completion status → attribute, not gate",
    "   (used to flag 'not assessed' individuals)",
    "",
    "• job_grade & assessment_level → feed the",
    "   Grade Fit score (#4)",
    "",
    "• year_in_sg → feeds Tenure Readiness (#5)",
    "",
    "• passed_tpcp → flag surfaced in output,",
    "   also boosts TPCP Health (#1)",
    "",
    "• skill_group / discipline / sub_discipline →",
    "   grouping dimensions for gap analysis",
]
add_bullets(slide, Inches(7.2), Inches(2.5), Inches(5.4), Inches(4.5), signals, size=12, color=DARK_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — New Pipeline Architecture
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "New Pipeline Architecture")

# Horizontal flow: Query → Join → Scope → 3 parallel tracks → Insights → Output
# Row 1 boxes
y1 = Inches(1.9)
w1 = Inches(2.4)
boxes = [
    ("HOD/GM Query\n(team + question)", PALE_BLUE, PRIMARY_BLUE),
    ("Employee Join\nHRIS ⋈ TPCP\non dummy_id", PALE_GREEN, GREEN),
    ("Org Scope Filter\nbusiness / opu /\ndivision", PALE_RED, RED),
]
for i, (lbl, c, b) in enumerate(boxes):
    x = Inches(0.5) + i * (w1 + Inches(0.6))
    add_shape(slide, x, y1, w1, Inches(1.3), c, b)
    add_text(slide, x + Inches(0.15), y1 + Inches(0.15), w1 - Inches(0.3), Inches(1),
             lbl, size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    if i < 2:
        add_text(slide, x + w1 + Inches(0.1), y1 + Inches(0.45), Inches(0.4), Inches(0.5),
                 "→", size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Downward arrow
add_text(slide, Inches(3.0), Inches(3.3), Inches(1), Inches(0.4), "↓",
         size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Row 2: 3 parallel tracks
y2 = Inches(3.8)
w2 = Inches(3.8)
tracks = [
    ("Per-Person Scoring\n6 health scores\nper employee", LIGHT_BLUE, PRIMARY_BLUE),
    ("Team Aggregate Rollup\npass rates, level distribution,\navg CBS, avg TPCP", PALE_PURPLE, PURPLE),
    ("Gap Matrix Analysis\ncompetency cell heatmap,\npeople vs structural", YELLOW_BG, AMBER),
]
for i, (lbl, c, b) in enumerate(tracks):
    x = Inches(0.5) + i * (w2 + Inches(0.25))
    add_shape(slide, x, y2, w2, Inches(1.5), c, b)
    add_text(slide, x + Inches(0.2), y2 + Inches(0.2), w2 - Inches(0.4), Inches(1.1),
             lbl, size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Downward arrow
add_text(slide, Inches(6.2), Inches(5.4), Inches(1), Inches(0.4), "↓",
         size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Row 3: Insight engine + output
y3 = Inches(5.9)
add_shape(slide, Inches(0.5), y3, Inches(5.8), Inches(1.2), PALE_GREEN, GREEN)
add_text(slide, Inches(0.7), y3 + Inches(0.15), Inches(5.4), Inches(0.9),
         "Insight Engine\nrule-based action triage (training / exposure / hiring, 6–12m horizon)",
         size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(6.35), y3 + Inches(0.45), Inches(0.4), Inches(0.5), "→",
         size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(6.9), y3, Inches(6), Inches(1.2), DARK_BLUE)
add_text(slide, Inches(7.1), y3 + Inches(0.15), Inches(5.6), Inches(0.9),
         "Pure JSON Output\nteam snapshot · flagged individuals · ranked gaps · actions\n→ passed to conversation layer for narrative",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — Output for each question type
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Output for Each Question Type (JSON Examples)")

# Left column — two example outputs
add_shape(slide, Inches(0.5), Inches(1.7), Inches(6.2), Inches(5.4), LIGHT_GRAY, RGBColor(0xCB, 0xD5, 0xE1))
add_text(slide, Inches(0.7), Inches(1.8), Inches(5.8), Inches(0.4),
         "Q: 'How did my talents perform in TPCP?'", size=13, bold=True, color=PRIMARY_BLUE)
json1 = """{
  "team_size": 45,
  "assessed": 32,
  "not_assessed": 13,
  "pass_rate": 0.84,
  "level_distribution": {
    "Staff": 20, "Principal": 8,
    "Custodian": 2, "Not TP": 2
  },
  "avg_cti_met_pct": 0.74,
  "weakest_dimension": "pacing"
}"""
tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(5.8), Inches(2.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = json1
p.font.size = Pt(11)
p.font.name = "Consolas"
p.font.color.rgb = DARK_BLUE

add_text(slide, Inches(0.7), Inches(4.6), Inches(5.8), Inches(0.4),
         "Q: 'Who is promotable in my division?'", size=13, bold=True, color=PRIMARY_BLUE)
json2 = """{
  "promotable": [
    {"employee_name": "...",
     "job_grade": "P4", "next_grade": "P5",
     "combined_score": 0.82,
     "year_in_sg": 4.2,
     "passed_tpcp": "Y",
     "gap_count": 2}
  ]
}"""
tb = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(5.8), Inches(2.1))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = json2
p.font.size = Pt(11)
p.font.name = "Consolas"
p.font.color.rgb = DARK_BLUE

# Right column
add_shape(slide, Inches(6.9), Inches(1.7), Inches(6.0), Inches(5.4), LIGHT_GRAY, RGBColor(0xCB, 0xD5, 0xE1))
add_text(slide, Inches(7.1), Inches(1.8), Inches(5.6), Inches(0.4),
         "Q: 'Who still has critical gaps?'", size=13, bold=True, color=PRIMARY_BLUE)
json3 = """{
  "flagged": [
    {"employee_name": "...",
     "qualified_level": "Staff",
     "combined_score": 0.42,
     "gap_count": 8,
     "weakest_cells": ["k3", "p2", "e1"]}
  ]
}"""
tb = slide.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.6), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = json3
p.font.size = Pt(11)
p.font.name = "Consolas"
p.font.color.rgb = DARK_BLUE

add_text(slide, Inches(7.1), Inches(4.1), Inches(5.6), Inches(0.4),
         "Q: 'What actions should we take?'", size=13, bold=True, color=PRIMARY_BLUE)
json4 = """{
  "actions": [
    {"type": "training", "priority": "high",
     "reason": "structural gap in Pacing",
     "affects": 18},
    {"type": "exposure", "priority": "medium",
     "reason": "Principals missing E3–E5",
     "affects": 4}
  ]
}"""
tb = slide.shapes.add_textbox(Inches(7.1), Inches(4.55), Inches(5.6), Inches(2.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = json4
p.font.size = Pt(11)
p.font.name = "Consolas"
p.font.color.rgb = DARK_BLUE


# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — Gap Analysis: People vs Structural
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Gap Analysis — People vs Structural",
           "Answers: 'Is this a people issue or a structural competency gap?'")

# Explanation
add_text(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
         "For each competency cell (b1–b25, k1–k15, p1–p10, e1–e10), compute the team failure rate:",
         size=14, color=DARK_BLUE)

headers = ["Failure Rate", "Classification", "Interpretation", "Recommended Action"]
rows = [
    ["≥ 60%", "STRUCTURAL", "Team-wide deficiency — not an individual problem",
     "Systemic training program, curriculum review"],
    ["20% – 60%", "MIXED", "Sub-group weakness — targeted pockets",
     "Targeted intervention, peer mentoring"],
    ["< 20%", "PEOPLE ISSUE", "Individual performance variance",
     "Individual coaching, performance management"],
]
add_table(slide, Inches(0.6), Inches(2.4),
          [Inches(1.8), Inches(2.2), Inches(4.2), Inches(4)],
          Inches(1), headers, rows, header_font=13, body_font=12)

# Repeat-gap detection
add_shape(slide, Inches(0.6), Inches(6.1), Inches(12.3), Inches(1), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(0.9), Inches(6.2), Inches(11.7), Inches(0.4),
         "Repeat-Gap Detection  (answers: 'Are there repeat gaps across multiple talents?')",
         size=13, bold=True, color=PRIMARY_BLUE)
add_text(slide, Inches(0.9), Inches(6.55), Inches(11.7), Inches(0.5),
         "Rank competency cells by count of employees scoring ≤ 1. Top N surfaces the systemic weak spots.",
         size=12, color=DARK_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — Action Recommendation Logic
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Action Recommendation Logic",
           "Rule-based triage across training / exposure / hiring")

# Three columns
cols = [
    {
        "title": "TRAINING",
        "color": LIGHT_BLUE, "border": PRIMARY_BLUE,
        "triggers": [
            "Triggered when:",
            "• Gap is STRUCTURAL (≥60% fail rate)",
            "• Weakness in Base (B) or Key (K) cells",
            "• Gap closable at level 2 → 3",
            "",
            "Horizon: 3–6 months",
            "Cost: Moderate (scale helps)",
            "Impact: High if structural",
        ],
    },
    {
        "title": "EXPOSURE",
        "color": YELLOW_BG, "border": AMBER,
        "triggers": [
            "Triggered when:",
            "• Gap in Pacing (P) or Emerging (E) cells",
            "• Competency at level 2 needing level 3",
            "• Individual issue, not structural",
            "",
            "Horizon: 6–12 months",
            "Cost: Low (stretch assignments)",
            "Impact: High for promotable talent",
        ],
    },
    {
        "title": "HIRING",
        "color": PALE_RED, "border": RED,
        "triggers": [
            "Triggered when:",
            "• Grade distribution skewed",
            "   (too few Principals for team)",
            "• Gap cells at level 0 across team",
            "   (no internal exposure at all)",
            "• Critical role unfilled",
            "",
            "Horizon: 12+ months",
            "Cost: High",
            "Impact: Fills structural vacancies",
        ],
    },
]
w = Inches(4.1)
gap_x = Inches(0.1)
for i, c in enumerate(cols):
    x = Inches(0.4) + i * (w + gap_x)
    y = Inches(1.75)
    add_shape(slide, x, y, w, Inches(5.3), c["color"], c["border"])
    add_text(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.5),
             c["title"], size=20, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_bullets(slide, x + Inches(0.25), y + Inches(0.8), w - Inches(0.5), Inches(4.3),
                c["triggers"], size=12, color=DARK_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 13 — Promotability Rules
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Promotability Logic",
           "Answers: 'Who is promotable?' · 'What gaps before promotion?'")

# Criteria
add_shape(slide, Inches(0.6), Inches(1.75), Inches(6), Inches(5.2), PALE_GREEN, GREEN)
add_text(slide, Inches(0.9), Inches(1.9), Inches(5.4), Inches(0.5),
         "Promotion-Ready (ALL conditions)", size=17, bold=True, color=GREEN)
rdy = [
    "✓  Combined Competency Score ≥ 0.75",
    "",
    "✓  year_in_sg ≥ grade_tenure_threshold",
    "    (P1–P3 → 3y, P4–P6 → 4y, P7+ → 5y)",
    "",
    "✓  passed_tpcp = 'Y'",
    "    (or no TPCP gate at current grade)",
    "",
    "✓  Critical Gap Count ≤ 3",
    "    (no more than 3 cells at level 0–1)",
    "",
    "✓  Grade Fit = 1.0 or 0.5",
    "    (not underqualified)",
]
add_bullets(slide, Inches(0.9), Inches(2.45), Inches(5.4), Inches(4.3), rdy, size=13)

# Gap-before-promotion
add_shape(slide, Inches(6.9), Inches(1.75), Inches(6), Inches(5.2), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(7.2), Inches(1.9), Inches(5.4), Inches(0.5),
         "Gaps Before Promotion", size=17, bold=True, color=PRIMARY_BLUE)
gaps = [
    "For employees scoring 0.65–0.74 (close to ready):",
    "",
    "• List competency cells at level 0 or 1",
    "• Flag any Key-level (K1–K15) gap as blocker",
    "• Flag passed_tpcp = 'N' as blocker",
    "• Show year_in_sg vs threshold delta",
    "",
    "Output includes 'close-in-6-months' flag",
    "for gaps at level 2 (one step away).",
    "",
    "For employees at 0.55–0.64 (develop further):",
    "• Return full gap profile, 12+ month plan",
]
add_bullets(slide, Inches(7.2), Inches(2.45), Inches(5.4), Inches(4.3), gaps, size=13)


# ════════════════════════════════════════════════════════════════════
# SLIDE 14 — New API Surface
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "New API Surface")

headers = ["Endpoint", "Purpose"]
rows = [
    ["POST /team/analyze",
     "Full diagnostic for a team: snapshot + individuals + gaps + actions. Body: {business, opu, division, question?}"],
    ["GET  /employee/{dummy_id}",
     "Individual scorecard: 6 health scores, gap profile, promotion readiness flags"],
    ["GET  /team/{business}/{opu}/{division}/promotable",
     "Ranked list of promotion-ready employees with rationale"],
    ["GET  /team/{business}/{opu}/{division}/gaps",
     "Competency gap matrix + people-vs-structural classification"],
    ["GET  /team/{business}/{opu}/{division}/snapshot",
     "Team aggregate: pass rate, level distribution, avg CBS/TPCP, not-assessed count"],
    ["POST /chat",
     "Natural language entry point — parses team + question type and routes to the right endpoint above"],
    ["GET  /health  ·  /docs  ·  /history",
     "Retained from v1 (liveness, Swagger UI, audit trail)"],
]
add_table(slide, Inches(0.4), Inches(1.75),
          [Inches(4.8), Inches(7.7)],
          Inches(0.65), headers, rows, header_font=13, body_font=12)


# ════════════════════════════════════════════════════════════════════
# SLIDE 15 — Module Structure & Implementation Order
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Module Structure & Implementation Order")

# Left — module layout
add_shape(slide, Inches(0.6), Inches(1.75), Inches(6.2), Inches(5.3), LIGHT_GRAY, RGBColor(0xCB, 0xD5, 0xE1))
add_text(slide, Inches(0.9), Inches(1.9), Inches(5.6), Inches(0.4),
         "app/ layout", size=15, bold=True, color=DARK_BLUE)
mod = """app/
├── data_loader.py     # load HRIS + TPCP, left-join
├── scores.py          # 6 health score functions
├── gap_analysis.py    # cell matrix, people vs structural
├── team_rollup.py     # team-level aggregate metrics
├── insight_engine.py  # action triage rules
├── recommender.py     # orchestrates full pipeline
├── config.py          # thresholds, weights
└── main.py            # new FastAPI endpoints"""
tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.35), Inches(5.8), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = mod
p.font.size = Pt(12)
p.font.name = "Consolas"
p.font.color.rgb = DARK_BLUE

# Right — implementation order
add_shape(slide, Inches(7.0), Inches(1.75), Inches(5.9), Inches(5.3), PALE_BLUE, PRIMARY_BLUE)
add_text(slide, Inches(7.3), Inches(1.9), Inches(5.3), Inches(0.4),
         "Implementation order", size=15, bold=True, color=PRIMARY_BLUE)
order = [
    "1.  data_loader.py  —  join HRIS + TPCP",
    "2.  scores.py  —  6 health score functions",
    "3.  team_rollup.py + gap_analysis.py",
    "4.  insight_engine.py  —  action triage",
    "5.  recommender.py + config.py  —  orchestrate",
    "6.  main.py  —  new FastAPI endpoints",
    "7.  chat parser  —  route NL to endpoints",
    "8.  frontend  —  dashboard view",
    "9.  evaluation  —  gap-detection metrics",
]
add_bullets(slide, Inches(7.3), Inches(2.45), Inches(5.3), Inches(4.3), order, size=14)


# ════════════════════════════════════════════════════════════════════
# SLIDE 16 — Migration Impact
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Migration Impact — What Changes")

headers = ["Component", "Status", "Action"]
rows = [
    ["data_loader.py", "REWRITE", "Load 2 JSON files, left-join on dummy_id"],
    ["features.py", "REPLACE", "Becomes scores.py — 6 absolute health functions"],
    ["recommender.py", "REWRITE", "New 4-stage pipeline: join → scope → score → insight"],
    ["config.py", "REWRITE", "New weights, grade-tenure thresholds, failure-rate cutoffs"],
    ["main.py (FastAPI)", "REWRITE", "/recommend retired, new team-centric endpoints"],
    ["history.py", "KEEP", "Audit trail still relevant"],
    ["chat/parser.py", "REWRITE", "Parse team scope + question type, not project fields"],
    ["chat/orchestrator.py", "REWRITE", "Route parsed query to right endpoint"],
    ["chat/formatter.py", "REWRITE", "Convert JSON → narrative insights (conversation layer)"],
    ["evaluation/", "REWRITE", "Gap-detection accuracy replaces precision@K as primary metric"],
    ["frontend/", "REWRITE", "Dashboard replaces form-search (team overview + drill-downs)"],
]
add_table(slide, Inches(0.4), Inches(1.7),
          [Inches(3.2), Inches(1.5), Inches(7.8)],
          Inches(0.44), headers, rows, header_font=13, body_font=11)


# ════════════════════════════════════════════════════════════════════
# SLIDE 17 — What Stays the Same
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "What Stays the Same")

items_left = [
    "• Rule-based heuristic approach",
    "   (deterministic, no ML training)",
    "",
    "• Explainable scoring",
    "   (every number traceable to a formula)",
    "",
    "• Configurable weights and thresholds",
    "   (domain experts tune without retraining)",
    "",
    "• FastAPI backend + CORS-enabled frontend",
    "",
    "• Auto-generated API docs at /docs",
]
items_right = [
    "• Audit trail of every request/response",
    "",
    "• Same deployment model",
    "   (uvicorn + single-file frontend)",
    "",
    "• Evaluation framework stays",
    "   (metrics adapted to new engine)",
    "",
    "• Chat entry point retained",
    "   (now parses team questions, not project specs)",
    "",
    "• Same production readiness path",
]

add_shape(slide, Inches(0.6), Inches(1.75), Inches(6), Inches(5.3), PALE_GREEN, GREEN)
add_bullets(slide, Inches(0.9), Inches(1.95), Inches(5.4), Inches(5), items_left, size=13, color=DARK_BLUE)

add_shape(slide, Inches(6.9), Inches(1.75), Inches(6), Inches(5.3), PALE_GREEN, GREEN)
add_bullets(slide, Inches(7.2), Inches(1.95), Inches(5.4), Inches(5), items_right, size=13, color=DARK_BLUE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 18 — Summary & Next Steps
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header(slide, "Summary & Next Steps")

# Summary
add_text(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
         "Summary", size=18, bold=True, color=PRIMARY_BLUE)
summary = [
    "• Engine repositioned from project matching to team diagnostics",
    "• Old 4 scoring functions retired — 6 absolute health scores replace them",
    "• Single mandatory hard filter: org scope (business / opu / division)",
    "• TPCP and CBS combined into a unified competency signal",
    "• Outputs pure JSON — narrative handled by the conversation layer",
    "• Answers all 16 HOD/GM questions across 4 analytical workflows",
]
add_bullets(slide, Inches(0.8), Inches(2.15), Inches(11.8), Inches(2.5), summary, size=14)

# Next steps
add_text(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
         "Decisions Needed Before Build", size=18, bold=True, color=AMBER)
next_steps = [
    "1.  Sign-off on the 6 health score formulas (weights are tunable, but the shape is fixed)",
    "2.  Confirm the action-triage thresholds (structural ≥ 60%, mixed 20–60%, people < 20%)",
    "3.  Confirm the promotability criteria (score ≥ 0.75, tenure, passed TPCP, ≤ 3 gaps)",
    "4.  Approval to retire the /recommend endpoint and replace with /team/* endpoints",
]
add_bullets(slide, Inches(0.8), Inches(5.25), Inches(11.8), Inches(2), next_steps, size=14)


# ════════════════════════════════════════════════════════════════════
# SLIDE 19 — Thank You
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "Ready to Build", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
         "Awaiting sign-off on scores, thresholds, and endpoint migration",
         size=18, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
         "Questions & Discussion",
         size=20, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
prs.save(str(OUTPUT_FILE))
print(f"Presentation saved to: {OUTPUT_FILE}")
print(f"Total slides: {len(prs.slides)}")
